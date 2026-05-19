# راوتر لوحة المدير
from fastapi import APIRouter, HTTPException, Depends, status, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from app.models.schemas import SchoolSetupRequest, StatsResponse
from app.services.account_generator import generate_accounts
from app.auth import get_current_user
from app.database import get_db
import io
import csv
import logging

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/manager", tags=["المدير"])


def require_manager(current_user: dict = Depends(get_current_user)):
    """التحقق من دور المدير"""
    if current_user["role"] not in ("manager", "admin", "owner"):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="هذه الصفحة للمديرين فقط"
        )
    return current_user


@router.get("/schools")
async def get_schools(
    current_user: dict = Depends(require_manager),
    db=Depends(get_db)
):
    """جلب قائمة المدارس"""
    result = db.table("schools").select("*").execute()
    return result.data


@router.post("/setup")
async def setup_school(
    request: SchoolSetupRequest,
    current_user: dict = Depends(require_manager),
    db=Depends(get_db)
):
    """إعداد المدرسة وتوليد الحسابات"""

    # التحقق من وجود المدرسة
    school = db.table("schools").select("*").eq("id", request.school_id).execute()
    if not school.data:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="المدرسة غير موجودة"
        )

    school_name = school.data[0]["name"]

    # توليد الحسابات
    accounts = await generate_accounts(
        school_id=request.school_id,
        students=request.students,
        teachers=request.teachers,
        admins_count=request.admins_count,
        db=db
    )

    # حفظ بيانات الإعداد مع الباسووردات
    db.table("school_setup").upsert({
        "school_id": request.school_id,
        "students_data": request.students,
        "teachers_data": request.teachers,
        "admins_count": request.admins_count,
        "total_accounts_generated": len(accounts),
        "passwords_data": accounts  # حفظ كل بيانات الحسابات بما فيها الباسوورد
    }).execute()

    # تحديث حالة المدرسة
    db.table("schools").update({"setup_completed": True}).eq("id", request.school_id).execute()

    return {
        "accounts": accounts,
        "total": len(accounts),
        "school_name": school_name
    }


@router.get("/accounts/{school_id}")
async def get_accounts(
    school_id: str,
    current_user: dict = Depends(require_manager),
    db=Depends(get_db)
):
    """جلب جميع حسابات مدرسة"""
    result = db.table("users") \
        .select("id, email, role, full_name, grade, subject, ministry_id, is_active, created_at") \
        .eq("school_id", school_id) \
        .neq("role", "manager") \
        .execute()

    return result.data


@router.get("/export/{school_id}")
async def export_accounts_csv(
    school_id: str,
    current_user: dict = Depends(require_manager),
    db=Depends(get_db)
):
    """تحميل الحسابات كملف CSV مع دعم العربي"""

    result = db.table("users") \
        .select("full_name, email, role, grade, subject, ministry_id") \
        .eq("school_id", school_id) \
        .neq("role", "manager") \
        .execute()

    # إنشاء CSV مع BOM للعربي
    output = io.StringIO()
    writer = csv.writer(output)

    writer.writerow(["الاسم الكامل", "الإيميل", "الدور", "الصف", "المادة", "الرقم الوزاري"])
    for user in result.data:
        role_ar = {"student": "طالب", "teacher": "معلم", "admin": "إداري"}.get(user["role"], user["role"])
        writer.writerow([
            user.get("full_name", ""),
            user.get("email", ""),
            role_ar,
            user.get("grade", ""),
            user.get("subject", ""),
            user.get("ministry_id", ""),
        ])

    csv_content = "﻿" + output.getvalue()  # BOM للعربي

    school = db.table("schools").select("name").eq("id", school_id).execute()
    school_name = school.data[0]["name"] if school.data else "school"

    return StreamingResponse(
        io.BytesIO(csv_content.encode("utf-8-sig")),
        media_type="text/csv; charset=utf-8",
        headers={"Content-Disposition": f'attachment; filename="morix_accounts_{school_name}.csv"'}
    )


@router.get("/stats")
async def get_stats(
    current_user: dict = Depends(require_manager),
    db=Depends(get_db)
):
    """لوحة الإحصائيات الشاملة"""
    from datetime import datetime, timedelta, timezone

    school_id = current_user.get("school_id")

    # ======= المستخدمون =======
    users_query = db.table("users").select("id, role, learning_style, created_at", count="exact")
    if school_id:
        users_query = users_query.eq("school_id", school_id)
    users = users_query.execute()

    total_users = len(users.data)
    students     = [u for u in users.data if u["role"] == "student"]
    teachers     = [u for u in users.data if u["role"] == "teacher"]
    admins       = [u for u in users.data if u["role"] == "admin"]
    student_ids  = [u["id"] for u in students]
    teacher_ids  = [u["id"] for u in teachers]

    # أساليب التعلم
    learning_styles = {"visual": 0, "auditory": 0, "kinesthetic": 0, "unknown": 0}
    for s in students:
        style = s.get("learning_style") or "unknown"
        learning_styles[style] = learning_styles.get(style, 0) + 1

    # ======= المحادثات =======
    total_conversations = 0
    convs_data = []
    try:
        convs_q = db.table("ai_conversations").select("id", count="exact")
        if school_id:
            # school_id قد لا يكون موجوداً بعد migration_v5 — نستخدم student_id كبديل
            try:
                convs_q = convs_q.eq("school_id", school_id)
                convs_r = convs_q.execute()
                total_conversations = convs_r.count or len(convs_r.data)
                convs_data = convs_r.data
            except Exception:
                # fallback: filter by student_ids
                if student_ids:
                    convs_q2 = db.table("ai_conversations").select("id", count="exact") \
                        .in_("student_id", student_ids)
                    convs_r2 = convs_q2.execute()
                    total_conversations = convs_r2.count or len(convs_r2.data)
                    convs_data = convs_r2.data
        else:
            convs_r = convs_q.execute()
            total_conversations = convs_r.count or len(convs_r.data)
            convs_data = convs_r.data
    except Exception:
        total_conversations = 0

    # ======= الرسائل =======
    msgs_total = 0
    try:
        if convs_data:
            conv_ids = [c["id"] for c in convs_data]
            if conv_ids:
                msgs_q = db.table("ai_messages").select("id", count="exact") \
                    .in_("conversation_id", conv_ids)
                msgs_total = msgs_q.execute().count or 0
    except Exception:
        msgs_total = 0

    # ======= الواجبات والاختبارات والأوراق =======
    hw_count, tests_count, ws_count, hw_sub_count, test_res_count = 0, 0, 0, 0, 0
    try:
        hw_q = db.table("homework").select("id", count="exact")
        if teacher_ids:
            hw_q = hw_q.in_("teacher_id", teacher_ids)
        elif school_id:
            pass  # can't filter without teacher_ids when no school boundary
        hw_count = hw_q.execute().count or 0

        tst_q = db.table("tests").select("id", count="exact")
        if teacher_ids:
            tst_q = tst_q.in_("teacher_id", teacher_ids)
        tests_count = tst_q.execute().count or 0

        ws_q = db.table("worksheets").select("id", count="exact")
        if teacher_ids:
            ws_q = ws_q.in_("teacher_id", teacher_ids)
        ws_count = ws_q.execute().count or 0

        if student_ids:
            hw_sub_count = db.table("homework_submissions").select("id", count="exact") \
                .in_("student_id", student_ids).execute().count or 0
            test_res_count = db.table("test_results").select("id", count="exact") \
                .in_("student_id", student_ids).execute().count or 0
    except Exception:
        pass

    # ======= الكتب =======
    books_count = 0
    try:
        # curriculum_books مشتركة بين كل المدارس (لا يوجد school_id حتى الآن)
        books_count = db.table("curriculum_books").select("id", count="exact").execute().count or 0
    except Exception:
        pass

    # ======= الألعاب والجلسات =======
    # focus_sessions مخزّنة في analytics بـ event_type='focus_session'
    games_count, focus_sessions_count = 0, 0
    try:
        if student_ids:
            games_count = db.table("educational_games").select("id", count="exact") \
                .in_("student_id", student_ids).execute().count or 0
        # جلسات التركيز من analytics (لا يوجد جدول منفصل)
        fs_q = db.table("analytics").select("id", count="exact").eq("event_type", "focus_session")
        if school_id:
            fs_q = fs_q.eq("school_id", school_id)
        focus_sessions_count = fs_q.execute().count or 0
    except Exception:
        pass

    # ======= الشارات والإنجازات =======
    # badges جدول يستخدم user_id (لا student_id)
    badges_earned = 0
    try:
        if student_ids:
            badges_earned = db.table("badges").select("id", count="exact") \
                .in_("user_id", student_ids).execute().count or 0
    except Exception:
        pass

    # ======= تسجيلات الدخول (7 أيام) =======
    now_utc = datetime.now(timezone.utc)
    week_ago = (now_utc - timedelta(days=7)).isoformat()
    month_ago = (now_utc - timedelta(days=30)).isoformat()

    logins_week, logins_month = 0, 0
    try:
        logins_week  = db.table("analytics").select("id", count="exact") \
            .eq("event_type", "login").gt("created_at", week_ago).execute().count or 0
        logins_month = db.table("analytics").select("id", count="exact") \
            .eq("event_type", "login").gt("created_at", month_ago).execute().count or 0
    except Exception:
        pass

    # ======= المستخدمون النشطون (آخر 7 أيام) =======
    active_users_week = 0
    try:
        # analytics يستخدم student_id (وإن كان المستخدم أي دور)
        active_q = db.table("analytics").select("student_id").gt("created_at", week_ago)
        if school_id:
            active_q = active_q.eq("school_id", school_id)
        active_evts = active_q.execute()
        active_users_week = len(set(e["student_id"] for e in active_evts.data if e.get("student_id")))
    except Exception:
        pass

    # ======= المستخدمون الجدد (آخر 30 يوم) =======
    new_users_month = 0
    try:
        new_u = db.table("users").select("id", count="exact").gt("created_at", month_ago)
        if school_id:
            new_u = new_u.eq("school_id", school_id)
        new_users_month = new_u.execute().count or 0
    except Exception:
        pass

    # ======= متوسط نتيجة الاختبارات =======
    avg_test_score = None
    try:
        if student_ids:
            scores_r = db.table("test_results").select("score") \
                .in_("student_id", student_ids).execute()
            scores = [r["score"] for r in scores_r.data if r.get("score") is not None]
            avg_test_score = round(sum(scores) / len(scores), 1) if scores else None
    except Exception:
        pass

    # ======= توزيع الصفوف/الدرجات =======
    grades_dist: dict = {}
    try:
        if student_ids:
            g_data = db.table("users").select("grade").in_("id", student_ids).execute()
            for row in g_data.data:
                gr = row.get("grade") or "غير محدد"
                grades_dist[gr] = grades_dist.get(gr, 0) + 1
    except Exception:
        pass

    # ======= الشكاوى =======
    complaints_count = 0
    try:
        cmp_q = db.table("complaints").select("id", count="exact")
        if school_id:
            cmp_q = cmp_q.eq("school_id", school_id)
        complaints_count = cmp_q.execute().count or 0
    except Exception:
        pass

    return {
        # مستخدمون
        "total_users":       total_users,
        "total_students":    len(students),
        "total_teachers":    len(teachers),
        "total_admins":      len(admins),
        "new_users_month":   new_users_month,
        "active_users_week": active_users_week,
        # ذكاء اصطناعي
        "total_conversations": total_conversations,
        "total_ai_messages":   msgs_total,
        # أكاديمي
        "total_homework":      hw_count,
        "total_tests":         tests_count,
        "total_worksheets":    ws_count,
        "homework_submissions":hw_sub_count,
        "test_results":        test_res_count,
        "avg_test_score":      avg_test_score,
        "total_books":         books_count,
        # نشاط
        "games_played":        games_count,
        "focus_sessions":      focus_sessions_count,
        "badges_earned":       badges_earned,
        "complaints_count":    complaints_count,
        # تسجيل دخول
        "logins_week":  logins_week,
        "logins_month": logins_month,
        # توزيعات
        "learning_styles": learning_styles,
        "grades_dist":     grades_dist,
        # توافق مع الإصدار القديم
        "recent_logins": logins_week,
    }


@router.get("/passwords/{school_id}")
async def get_saved_passwords(
    school_id: str,
    current_user: dict = Depends(require_manager),
    db=Depends(get_db)
):
    """جلب الباسووردات المحفوظة للمدير"""
    result = db.table("school_setup") \
        .select("passwords_data, total_accounts_generated, updated_at") \
        .eq("school_id", school_id) \
        .execute()

    if not result.data or not result.data[0].get("passwords_data"):
        return {"accounts": [], "total": 0}

    return {
        "accounts": result.data[0]["passwords_data"],
        "total": result.data[0]["total_accounts_generated"],
        "generated_at": result.data[0]["updated_at"]
    }


@router.get("/books")
async def get_books(
    current_user: dict = Depends(require_manager),
    db=Depends(get_db)
):
    """جلب كتب المنهج"""
    result = db.table("curriculum_books").select("*").eq("is_active", True).execute()
    return result.data


@router.post("/extract-book-text")
async def extract_book_text(
    file: UploadFile = File(...),
    current_user: dict = Depends(require_manager),
):
    """استخراج النص من ملف PDF أو PowerPoint أو TXT"""
    filename = (file.filename or "").lower()
    content = await file.read()

    if len(content) > 15 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="الملف أكبر من 15MB")

    text = ""

    # ───── TXT / MD ─────
    if filename.endswith((".txt", ".md")):
        try:
            text = content.decode("utf-8", errors="replace")
        except Exception:
            text = content.decode("latin-1", errors="replace")

    # ───── PDF ─────
    elif filename.endswith(".pdf"):
        try:
            from pypdf import PdfReader
            import io as _io
            reader = PdfReader(_io.BytesIO(content))
            pages_text = []
            for page in reader.pages:
                t = page.extract_text() or ""
                pages_text.append(t)
            text = "\n".join(pages_text)
        except Exception as e:
            logger.warning(f"pypdf failed: {e}")
            # fallback: raw decode (text-layer PDFs)
            try:
                text = content.decode("utf-8", errors="ignore")
            except Exception:
                text = ""

    # ───── PPTX ─────
    elif filename.endswith(".pptx"):
        try:
            from pptx import Presentation
            import io as _io
            prs = Presentation(_io.BytesIO(content))
            slides_text = []
            for slide in prs.slides:
                slide_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                if slide_parts:
                    slides_text.append(" | ".join(slide_parts))
            text = "\n".join(slides_text)
        except Exception as e:
            logger.warning(f"pptx extraction failed: {e}")
            text = ""

    else:
        raise HTTPException(status_code=400, detail="صيغة غير مدعومة. استخدم PDF أو PPTX أو TXT")

    if not text.strip():
        raise HTTPException(status_code=422, detail="لم يتم استخراج أي نص من الملف — تأكد أن الملف يحتوي على نص قابل للقراءة")

    # تقليص النص لـ 80,000 حرف (حد Gemini المناسب)
    text = text[:80000]
    return {"text": text, "chars": len(text), "filename": file.filename}


@router.post("/books")
async def add_book(
    book_data: dict,
    current_user: dict = Depends(require_manager),
    db=Depends(get_db)
):
    """إضافة كتاب منهجي جديد"""
    from app.services.book_summarizer import summarize_book

    title = book_data.get("title", "")
    subject = book_data.get("subject", "")
    raw_text = book_data.get("raw_text", "")

    summary = None
    if raw_text:
        summary = await summarize_book(title, subject, raw_text)

    result = db.table("curriculum_books").insert({
        "title": title,
        "subject": subject,
        "grade": book_data.get("grade", ""),
        "summary": summary,
        "key_concepts": book_data.get("key_concepts", []),
        "is_active": True
    }).execute()

    return result.data[0] if result.data else {"message": "تم إضافة الكتاب"}


# ============================================================
# 👔 ميزات المدير المتقدمة
# ============================================================

@router.post("/schools")
async def create_school(
    body: dict,
    current_user: dict = Depends(require_manager),
    db=Depends(get_db),
):
    """إضافة مدرسة جديدة (تنحفظ بشكل دائم في قاعدة البيانات)"""
    import uuid as _uuid
    name = (body.get("name") or "").strip()
    branch = (body.get("branch") or "").strip()
    ministry_code = (body.get("ministry_code") or "").strip()

    if not name:
        raise HTTPException(400, "اكتب اسم المدرسة")

    full_name = f"{name} - فرع {branch}" if branch else name

    # تأكد إن المدرسة مش مكررة
    existing = db.table("schools").select("id").eq("name", full_name).execute()
    if existing.data:
        raise HTTPException(400, "مدرسة بنفس الاسم والفرع موجودة بالفعل")

    # توليد كود وزاري تلقائي لو فاضي (الـ DB قد تكون بـ UNIQUE NOT NULL)
    if not ministry_code:
        ministry_code = f"MX-{_uuid.uuid4().hex[:8].upper()}"

    # محاولات تدريجية: مع كل الأعمدة، بدون branch، بدون كل الإضافات
    attempts = [
        {"name": full_name, "ministry_code": ministry_code, "branch": branch or None, "setup_completed": False},
        {"name": full_name, "ministry_code": ministry_code, "setup_completed": False},
        {"name": full_name, "ministry_code": ministry_code},
    ]

    last_error = None
    for payload in attempts:
        try:
            result = db.table("schools").insert(payload).execute()
            if result.data:
                return result.data[0]
        except Exception as e:
            last_error = e
            continue

    raise HTTPException(500, f"تعذرت إضافة المدرسة: {last_error}")


@router.delete("/schools/{school_id}")
async def delete_school(
    school_id: str,
    current_user: dict = Depends(require_manager),
    db=Depends(get_db),
):
    """حذف مدرسة + كل حساباتها نهائياً"""
    try:
        db.table("users").delete().eq("school_id", school_id).execute()
    except Exception:
        pass
    try:
        db.table("school_setup").delete().eq("school_id", school_id).execute()
    except Exception:
        pass
    db.table("schools").delete().eq("id", school_id).execute()
    return {"message": "✅ تم حذف المدرسة وكل حساباتها"}


@router.post("/upload-excel")
async def upload_excel_setup(
    school_id: str = Form(...),
    file: UploadFile = File(...),
    current_user: dict = Depends(require_manager),
    db=Depends(get_db),
):
    """
    رفع ملف Excel لإعداد المدرسة.
    الأعمدة المتوقعة (مرنة):
      - الاسم / Name / full_name
      - الدور / Role / role  (student / teacher / admin / طالب / معلم / إداري)
      - الصف / Grade / grade  (للطالب)
      - المادة / Subject / subject  (للمعلم)
      - الرقم الوزاري / Ministry ID / ministry_id  (اختياري)
    """
    from openpyxl import load_workbook
    from app.services.account_generator import generate_accounts

    school = db.table("schools").select("*").eq("id", school_id).execute()
    if not school.data:
        raise HTTPException(404, "المدرسة غير موجودة")

    content = await file.read()
    try:
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        ws = wb.active
    except Exception as e:
        raise HTTPException(400, f"تعذرت قراءة ملف Excel: {e}")

    rows = list(ws.iter_rows(values_only=True))
    if not rows or len(rows) < 2:
        raise HTTPException(400, "الملف فارغ")

    headers = [str(h or "").strip().lower() for h in rows[0]]

    def find_col(*candidates):
        for c in candidates:
            for i, h in enumerate(headers):
                if c in h:
                    return i
        return -1

    col_name = find_col("اسم", "name", "full_name")
    col_role = find_col("دور", "role")
    col_grade = find_col("صف", "grade", "صفّ")
    col_subject = find_col("مادة", "subject")
    col_ministry = find_col("وزاري", "ministry", "id")

    if col_name == -1 or col_role == -1:
        raise HTTPException(400, "الملف لازم يحتوي على عمودين على الأقل: الاسم والدور")

    students, teachers, admins_count = [], [], 0
    role_map = {
        "طالب": "student", "student": "student",
        "معلم": "teacher", "معلّم": "teacher", "teacher": "teacher",
        "اداري": "admin", "إداري": "admin", "admin": "admin",
    }

    for row in rows[1:]:
        if not row or all(c is None for c in row):
            continue
        full_name = str(row[col_name] or "").strip()
        if not full_name:
            continue
        raw_role = str(row[col_role] or "").strip().lower()
        role = role_map.get(raw_role, "student")
        grade = str(row[col_grade] or "").strip() if col_grade != -1 else ""
        subject = str(row[col_subject] or "").strip() if col_subject != -1 else ""
        ministry_id = str(row[col_ministry] or "").strip() if col_ministry != -1 else ""

        entry = {"full_name": full_name}
        if ministry_id:
            entry["ministry_id"] = ministry_id

        if role == "student":
            entry["grade"] = grade or "غير محدد"
            students.append(entry)
        elif role == "teacher":
            entry["subject"] = subject or "عام"
            teachers.append(entry)
        elif role == "admin":
            admins_count += 1

    accounts = await generate_accounts(
        school_id=school_id,
        students=students,
        teachers=teachers,
        admins_count=admins_count,
        db=db,
    )

    db.table("school_setup").upsert({
        "school_id": school_id,
        "students_data": students,
        "teachers_data": teachers,
        "admins_count": admins_count,
        "total_accounts_generated": len(accounts),
        "passwords_data": accounts,
    }).execute()

    db.table("schools").update({"setup_completed": True}).eq("id", school_id).execute()

    return {
        "accounts": accounts,
        "total": len(accounts),
        "school_name": school.data[0]["name"],
        "students_count": len(students),
        "teachers_count": len(teachers),
        "admins_count": admins_count,
    }


@router.get("/account-password/{user_id}")
async def get_account_password(
    user_id: str,
    current_user: dict = Depends(require_manager),
    db=Depends(get_db),
):
    """عرض كلمة سر حساب معين (للمدير فقط)"""
    user = db.table("users").select("school_id, email, full_name").eq("id", user_id).execute()
    if not user.data:
        return {
            "found": False,
            "message": "❌ الحساب غير موجود",
            "password": None,
        }

    school_id = user.data[0]["school_id"]
    email = user.data[0]["email"]
    full_name = user.data[0]["full_name"]

    try:
        setup = db.table("school_setup").select("passwords_data").eq("school_id", school_id).execute()
        if setup.data and setup.data[0].get("passwords_data"):
            for acc in setup.data[0]["passwords_data"]:
                if acc.get("email") == email:
                    return {
                        "found": True,
                        "email": email,
                        "password": acc.get("password"),
                        "full_name": full_name,
                    }
    except Exception as e:
        logger.warning(f"password lookup failed: {e}")

    # كلمة السر غير محفوظة → اقترح إعادة تعيين
    return {
        "found": False,
        "email": email,
        "full_name": full_name,
        "password": None,
        "message": "🔐 كلمة السر غير محفوظة لهذا الحساب (تم إنشاؤه قبل تفعيل ميزة الحفظ التلقائي).",
        "suggestion": "اضغط على 'إعادة تعيين' في صفحة الحسابات لإنشاء كلمة سر جديدة.",
    }


@router.put("/accounts/{user_id}/reset-password")
async def manager_reset_password(
    user_id: str,
    current_user: dict = Depends(require_manager),
    db=Depends(get_db),
):
    """إعادة تعيين كلمة سر حساب — يولّد جديدة ويحفظها في school_setup"""
    import secrets, string
    from app.auth import hash_password

    user = db.table("users").select("school_id, email, full_name, role").eq("id", user_id).execute()
    if not user.data:
        raise HTTPException(404, "الحساب غير موجود")
    if user.data[0].get("role") in ("manager", "owner"):
        raise HTTPException(403, "لا يمكن إعادة تعيين كلمة سر مدير أو مالك من هنا")

    # توليد كلمة سر قوية
    alphabet = string.ascii_letters + string.digits
    new_password = ''.join(secrets.choice(alphabet) for _ in range(10))

    # تحديث الـ hash
    db.table("users").update({
        "password_hash": hash_password(new_password),
        "must_change_password": False,
    }).eq("id", user_id).execute()

    # حفظها في school_setup
    school_id = user.data[0]["school_id"]
    email = user.data[0]["email"]
    full_name = user.data[0]["full_name"]
    role = user.data[0]["role"]
    try:
        setup = db.table("school_setup").select("passwords_data").eq("school_id", school_id).execute()
        if setup.data:
            pw_list = setup.data[0].get("passwords_data") or []
            updated = False
            for acc in pw_list:
                if acc.get("email") == email:
                    acc["password"] = new_password
                    updated = True
                    break
            if not updated:
                pw_list.append({"email": email, "password": new_password, "full_name": full_name, "role": role})
            db.table("school_setup").update({"passwords_data": pw_list}).eq("school_id", school_id).execute()
        else:
            db.table("school_setup").insert({
                "school_id": school_id,
                "passwords_data": [{"email": email, "password": new_password, "full_name": full_name, "role": role}],
            }).execute()
    except Exception as e:
        logger.warning(f"saving reset password failed: {e}")

    return {
        "message": f"✅ تم إعادة تعيين كلمة السر للحساب {email}",
        "email": email,
        "password": new_password,
        "full_name": full_name,
    }


@router.delete("/accounts/{user_id}")
async def delete_account(
    user_id: str,
    current_user: dict = Depends(require_manager),
    db=Depends(get_db),
):
    """حذف حساب نهائياً من قاعدة البيانات (مع كل بياناته)"""
    user = db.table("users").select("id, school_id, email, role").eq("id", user_id).execute()
    if not user.data:
        raise HTTPException(404, "الحساب غير موجود")

    if user.data[0].get("role") in ("manager", "owner"):
        raise HTTPException(403, "لا يمكن حذف حسابات المديرين أو الملاك من هنا")

    email = user.data[0]["email"]
    school_id = user.data[0]["school_id"]

    # حذف كل البيانات المرتبطة بالحساب
    related_tables = [
        "ai_conversations", "ai_messages", "messages", "conversations",
        "homework", "tests", "worksheets", "submissions",
        "diagnostic_results", "streaks", "achievements",
        "user_settings", "analytics", "complaints",
    ]
    for tbl in related_tables:
        try:
            db.table(tbl).delete().eq("user_id", user_id).execute()
        except Exception:
            pass
        try:
            db.table(tbl).delete().eq("student_id", user_id).execute()
        except Exception:
            pass

    # حذف الحساب من users
    db.table("users").delete().eq("id", user_id).execute()

    # حذف من school_setup.passwords_data
    try:
        setup = db.table("school_setup").select("passwords_data").eq("school_id", school_id).execute()
        if setup.data and setup.data[0].get("passwords_data"):
            new_pw = [a for a in setup.data[0]["passwords_data"] if a.get("email") != email]
            db.table("school_setup").update({"passwords_data": new_pw}).eq("school_id", school_id).execute()
    except Exception:
        pass

    return {"message": f"✅ تم حذف الحساب {email} نهائياً مع كل بياناته"}


# ============================================================
# ⚙️ إعدادات المدير
# ============================================================
@router.get("/settings")
async def get_manager_settings(current_user: dict = Depends(require_manager), db=Depends(get_db)):
    result = db.table("user_settings").select("*").eq("user_id", current_user["id"]).execute()
    base = {
        "theme": "dark", "language": "ar",
        "notifications_enabled": True,
        "avatar_url": current_user.get("avatar_url", ""),
        "email": current_user.get("email", ""),
        "full_name": current_user.get("full_name", ""),
    }
    if not result.data: return base
    s = result.data[0]
    return {**base, **{k: v for k, v in s.items() if k not in ("id","user_id","created_at","updated_at")}, "avatar_url": current_user.get("avatar_url", "")}


@router.put("/settings")
async def update_manager_settings(body: dict, current_user: dict = Depends(require_manager), db=Depends(get_db)):
    valid_langs = {"ar", "en", "de", "fr", "zh", "es"}
    valid_themes = {"dark", "light", "library", "neon"}
    theme = body.get("theme", "dark")
    if theme not in valid_themes: theme = "dark"
    lang = body.get("language", "ar")
    if lang not in valid_langs: lang = "ar"
    data = {
        "user_id": current_user["id"], "theme": theme,
        "language": lang, "notifications_enabled": bool(body.get("notifications_enabled", True)),
    }
    try:
        existing = db.table("user_settings").select("id").eq("user_id", current_user["id"]).execute()
        if existing.data:
            db.table("user_settings").update(data).eq("user_id", current_user["id"]).execute()
        else:
            db.table("user_settings").insert(data).execute()
    except Exception as e:
        logger.warning(f"Manager settings save failed: {e}")
    if body.get("avatar_url") is not None:
        try: db.table("users").update({"avatar_url": body["avatar_url"]}).eq("id", current_user["id"]).execute()
        except: pass
    return {"message": "✅ تم حفظ الإعدادات"}


@router.post("/strategic-advisor")
async def strategic_advisor(
    body: dict,
    current_user: dict = Depends(require_manager),
):
    """مستشار استراتيجي ذكي"""
    question = (body.get("question") or "").strip()[:2000]
    context = str(body.get("context", ""))[:500]
    if not question:
        raise HTTPException(400, "اطرح سؤالك الاستراتيجي")
    from app.services.ai_service import chat_with_gemini
    prompt = f"""أنت مستشار استراتيجي تعليمي خبير.
سياق المستخدم: {context or 'مدير مدارس متعددة'}
السؤال: {question}

قدم استشارة استراتيجية شاملة:
1. تحليل الموقف
2. 3 خيارات استراتيجية مع إيجابيات/سلبيات
3. التوصية النهائية مع خطة تنفيذ مرحلية
4. مؤشرات نجاح قابلة للقياس (KPIs)"""
    try:
        text, _ = await chat_with_gemini(prompt, None, full_name=current_user.get("full_name", ""), role="manager")
        return {"advice": text}
    except Exception:
        return {"advice": "تعذر الاتصال بـ AI"}


@router.get("/health-score")
async def franchise_health(current_user: dict = Depends(require_manager), db=Depends(get_db)):
    """درجة صحية شاملة لكل المدارس"""
    from datetime import datetime, timedelta, timezone
    week = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat()
    scored = []
    try:
        schools = db.table("schools").select("id, name, setup_completed").execute()
        for sc in (schools.data or []):
            try:
                users = db.table("users").select("id").eq("school_id", sc["id"]).execute()
                total = len(users.data or [])
            except Exception:
                total = 0

            if total == 0:
                scored.append({"school_id": sc["id"], "name": sc["name"], "score": 20, "status": "🔴 فارغة", "users": 0, "active_week": 0})
                continue

            active = 0
            try:
                ull = db.table("users").select("last_login").eq("school_id", sc["id"]).execute()
                active = sum(1 for u in (ull.data or []) if u.get("last_login") and str(u["last_login"]) >= week)
            except Exception:
                pass
            if active == 0:
                try:
                    logs = db.table("analytics").select("student_id").eq("school_id", sc["id"]).eq("event_type", "login").gte("created_at", week).execute()
                    active = len(set(l.get("student_id") for l in (logs.data or []) if l.get("student_id")))
                except Exception:
                    pass

            engagement = (active / total) * 100
            score = int(engagement * 0.6 + (40 if sc.get("setup_completed") else 0))
            status_label = "🟢 ممتازة" if score >= 75 else "🟡 جيدة" if score >= 50 else "🟠 ضعيفة" if score >= 25 else "🔴 حرجة"
            scored.append({"school_id": sc["id"], "name": sc["name"], "score": score, "status": status_label, "users": total, "active_week": active})
    except Exception as e:
        logger.error(f"Health failed: {e}")
    return sorted(scored, key=lambda x: -x["score"])
