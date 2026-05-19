# راوتر الذكاء الاصطناعي - Morix Platform
from fastapi import APIRouter, HTTPException, Depends, UploadFile, File
from app.models.schemas import ChatMessage, ChatResponse, ImageGenerateRequest
from app.services.ai_service import chat_with_gemini, generate_educational_image
from app.auth import get_current_user
from app.database import get_db
import logging

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/ai", tags=["الذكاء الاصطناعي"])


@router.post("/chat", response_model=ChatResponse)
async def chat(
    message: ChatMessage,
    current_user: dict = Depends(get_current_user),
    db=Depends(get_db)
):
    """الشات مع الذكاء الاصطناعي - يدعم الطالب والمعلم والمالك"""
    role = current_user["role"]

    conversation_id = message.conversation_id
    if not conversation_id:
        conv_data = {
            "student_id": current_user["id"],
            "school_id": current_user.get("school_id"),
            "title": message.message[:50] + ("..." if len(message.message) > 50 else ""),
        }
        if message.book_id:
            conv_data["book_id"] = message.book_id

        conv = db.table("ai_conversations").insert(conv_data).execute()
        conversation_id = conv.data[0]["id"]
    elif message.book_id:
        db.table("ai_conversations").update({"book_id": message.book_id}) \
            .eq("id", conversation_id).execute()

    history = db.table("ai_messages") \
        .select("role, content") \
        .eq("conversation_id", conversation_id) \
        .order("created_at") \
        .execute()

    book_summary = None
    if message.book_id:
        book = db.table("curriculum_books").select("summary, title").eq("id", message.book_id).execute()
        if book.data and book.data[0].get("summary"):
            book_summary = f"Book: {book.data[0]['title']}\n{book.data[0]['summary']}"

    # جلب إعدادات المستخدم (الصعوبة + اللغة المفضلة)
    difficulty = "medium"
    preferred_lang = None
    try:
        s = db.table("user_settings").select("difficulty, language").eq("user_id", current_user["id"]).execute()
        if s.data:
            difficulty = s.data[0].get("difficulty", "medium")
            # إذا حدد المستخدم لغة معيّنة نأخذها — وإلا نستخدم لغة الرسالة
            preferred_lang = s.data[0].get("language") or None
    except Exception:
        pass

    # لغة الرسالة تأخذ أولوية على إعدادات DB (إذا أرسلها الفرونت)
    if message.language and message.language != "ar":
        preferred_lang = message.language

    learning_style = current_user.get("learning_style") if role == "student" else None
    full_name = current_user.get("full_name", "مستخدم")

    # حفظ رسالة المستخدم أولاً (قبل نداء AI) حتى لا تضيع عند الخطأ
    db.table("ai_messages").insert({"conversation_id": conversation_id, "role": "user", "content": message.message}).execute()

    reply, from_cache = await chat_with_gemini(
        message=message.message,
        learning_style=learning_style,
        book_summary=book_summary,
        conversation_history=history.data if history.data else [],
        full_name=full_name,
        role=role,
        difficulty=difficulty,
        image_base64=message.image_base64,
        file_text=message.file_text,
        preferred_lang=preferred_lang,
    )

    db.table("ai_messages").insert({"conversation_id": conversation_id, "role": "assistant", "content": reply}).execute()

    try:
        db.table("analytics").insert({
            "student_id": current_user["id"],
            "school_id": current_user.get("school_id"),
            "event_type": "ai_chat",
            "event_data": {"from_cache": from_cache, "role": role}
        }).execute()
    except Exception:
        pass

    return ChatResponse(conversation_id=conversation_id, reply=reply, from_cache=from_cache)


@router.get("/conversations")
async def get_conversations(current_user: dict = Depends(get_current_user), db=Depends(get_db)):
    result = db.table("ai_conversations") \
        .select("id, title, created_at, book_id, curriculum_books(title)") \
        .eq("student_id", current_user["id"]) \
        .order("created_at", desc=True).execute()
    return result.data


@router.get("/conversations/{conversation_id}/messages")
async def get_messages(conversation_id: str, current_user: dict = Depends(get_current_user), db=Depends(get_db)):
    conv = db.table("ai_conversations") \
        .select("*").eq("id", conversation_id).eq("student_id", current_user["id"]).execute()
    if not conv.data:
        raise HTTPException(status_code=404, detail="المحادثة غير موجودة")

    messages = db.table("ai_messages") \
        .select("*").eq("conversation_id", conversation_id).order("created_at").execute()
    return {"conversation": conv.data[0], "messages": messages.data}


@router.delete("/conversations/{conversation_id}")
async def delete_conversation(conversation_id: str, current_user: dict = Depends(get_current_user), db=Depends(get_db)):
    conv = db.table("ai_conversations") \
        .select("id").eq("id", conversation_id).eq("student_id", current_user["id"]).execute()
    if not conv.data:
        raise HTTPException(status_code=404, detail="المحادثة غير موجودة")
    db.table("ai_conversations").delete().eq("id", conversation_id).execute()
    return {"message": "تم حذف المحادثة"}


@router.get("/books")
async def get_books_for_chat(current_user: dict = Depends(get_current_user), db=Depends(get_db)):
    result = db.table("curriculum_books").select("id, title, subject, grade, summary").eq("is_active", True).execute()
    return result.data


@router.post("/generate-image")
async def generate_image(request: ImageGenerateRequest, current_user: dict = Depends(get_current_user)):
    """توليد صورة تعليمية بالذكاء الاصطناعي"""
    if not request.prompt.strip():
        raise HTTPException(status_code=400, detail="الوصف فارغ")

    result = await generate_educational_image(request.prompt)

    if not result.get("success"):
        return {
            "success": False,
            "message": result.get("error") or "تعذر توليد الصورة",
            "image": None,
            "details": result.get("details", ""),
        }
    image_b64 = result["image"]

    return {"success": True, "image": image_b64, "prompt": request.prompt}


# ============================================================
# 📄 استخراج نص من ملف — متاح لجميع المستخدمين المسجلين
# ============================================================
@router.post("/extract-file")
async def extract_file_for_chat(
    file: UploadFile = File(...),
    current_user: dict = Depends(get_current_user),
):
    """استخراج النص من ملف PDF أو PPTX أو DOCX أو TXT لاستخدامه في الشات"""
    filename = (file.filename or "").lower()
    content = await file.read()
    if len(content) > 15 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="الملف أكبر من 15MB — اختر ملفاً أصغر")
    text = ""

    if filename.endswith((".txt", ".md")):
        try:
            text = content.decode("utf-8", errors="replace")
        except Exception:
            text = content.decode("latin-1", errors="replace")

    elif filename.endswith(".pdf"):
        try:
            from pypdf import PdfReader
            import io as _io
            reader = PdfReader(_io.BytesIO(content))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception as e:
            logger.warning(f"pypdf failed: {e}")
            text = content.decode("utf-8", errors="ignore")

    elif filename.endswith(".pptx"):
        try:
            from pptx import Presentation
            import io as _io
            prs = Presentation(_io.BytesIO(content))
            slides_text = []
            for slide in prs.slides:
                parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if parts:
                    slides_text.append(" | ".join(parts))
            text = "\n".join(slides_text)
        except Exception as e:
            logger.warning(f"pptx extraction failed: {e}")
            text = ""

    elif filename.endswith(".docx"):
        try:
            import zipfile, re as _re, io as _io
            with zipfile.ZipFile(_io.BytesIO(content)) as zf:
                with zf.open("word/document.xml") as doc_xml:
                    xml_content = doc_xml.read().decode("utf-8", errors="replace")
            text = _re.sub(r"<[^>]+>", " ", xml_content)
            text = _re.sub(r"\s+", " ", text).strip()
        except Exception as e:
            logger.warning(f"docx extraction failed: {e}")
            text = ""

    else:
        raise HTTPException(status_code=400, detail="صيغة غير مدعومة — استخدم PDF أو PPTX أو DOCX أو TXT أو MD")

    if not text.strip():
        raise HTTPException(status_code=422, detail="لم يُستخرج أي نص من الملف — تأكد أنه يحتوي على نص قابل للقراءة")

    text = text[:80000]
    return {"text": text, "chars": len(text), "filename": file.filename}
